import pandas as pd
import sqlite3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

conn = sqlite3.connect('financial_analytics.db')

print("=" * 70)
print("GENERATING POWERPOINT REPORT")
print("=" * 70)

# >>> FETCH DATA FOR SLIDES
query_summary = """
SELECT
    ROUND(SUM(revenue), 2) AS total_revenue,
    ROUND(AVG((gross_profit / revenue * 100)), 1) AS gross_margin_pct,
    ROUND(AVG((operating_income / revenue * 100)), 1) AS op_margin_pct,
    COUNT(DISTINCT customer_id) AS total_customers
FROM financial_data
"""

summary = pd.read_sql_query(query_summary, conn).iloc[0]

query_by_segment = """
SELECT
    customer_segment,
    ROUND(SUM(revenue), 2) AS revenue,
    ROUND(AVG((gross_profit / revenue * 100)), 1) AS margin_pct
FROM financial_data
GROUP BY customer_segment
ORDER BY revenue DESC
"""

by_segment = pd.read_sql_query(query_by_segment, conn)

query_by_product = """
SELECT
    product,
    ROUND(SUM(revenue), 2) AS revenue
FROM financial_data
GROUP BY product
ORDER BY revenue DESC
LIMIT 5
"""

by_product = pd.read_sql_query(query_by_product, conn)

# >>> CREATE PRESENTATION
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.8), width, Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    return slide

def add_metrics_slide(prs, title, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    
    left_start = Inches(1)
    top_start = Inches(1.8)
    box_width = Inches(4)
    box_height = Inches(1.5)
    
    for idx, (label, value) in enumerate(metrics):
        left = left_start + (Inches(4.5) if idx % 2 == 1 else Inches(0))
        top = top_start + (Inches(2) * (idx // 2))
        
        box = slide.shapes.add_shape(1, left, top, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(230, 240, 250)
        box.line.color.rgb = RGBColor(0, 102, 204)
        
        text_frame = box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100)
        
        p = text_frame.add_paragraph()
        p.text = str(value)
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

# >>> SLIDE 1: TITLE
print("Creating Slide 1: Title")
add_title_slide(prs, "Financial Dashboard", "2023-2024 Financial Analysis")

# >>> SLIDE 2: KEY METRICS
print("Creating Slide 2: Key Metrics")
metrics = [
    ("Total Revenue", f"${summary['total_revenue']:,.0f}"),
    ("Gross Margin", f"{summary['gross_margin_pct']:.1f}%"),
    ("Operating Margin", f"{summary['op_margin_pct']:.1f}%"),
    ("Total Customers", f"{int(summary['total_customers']):,}")
]
add_metrics_slide(prs, "Financial Performance", metrics)

# >>> SLIDE 3: REVENUE BY SEGMENT
print("Creating Slide 3: Revenue by Segment")
slide = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Revenue by Customer Segment"
p.font.size = Pt(44)
p.font.bold = True

rows = len(by_segment) + 1
cols = 3

left, top, width, height = Inches(1), Inches(1.8), Inches(8), Inches(4)
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

headers = ["Segment", "Revenue", "Margin %"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True

for row_idx, row_data in enumerate(by_segment.itertuples(index=False), 1):
    table.cell(row_idx, 0).text = str(row_data[0])
    table.cell(row_idx, 1).text = f"${row_data[1]:,.0f}"
    table.cell(row_idx, 2).text = f"{row_data[2]:.1f}%"

# >>> SLIDE 4: TOP PRODUCTS
print("Creating Slide 4: Top Products")
slide = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Top 5 Products by Revenue"
p.font.size = Pt(44)
p.font.bold = True

rows = len(by_product) + 1
cols = 2

left, top, width, height = Inches(2), Inches(1.8), Inches(6), Inches(4)
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

headers = ["Product", "Revenue"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True

for row_idx, row_data in enumerate(by_product.itertuples(index=False), 1):
    table.cell(row_idx, 0).text = str(row_data[0])
    table.cell(row_idx, 1).text = f"${row_data[1]:,.0f}"

# >>> SLIDE 5: SUMMARY
print("Creating Slide 5: Summary")
slide = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Key Takeaways"
p.font.size = Pt(44)
p.font.bold = True

takeaways = [
    f"Total 24-month revenue: ${summary['total_revenue']:,.0f}",
    f"Healthy gross margins averaging {summary['gross_margin_pct']:.1f}%",
    f"Operating margins at {summary['op_margin_pct']:.1f}% indicate strong cost control",
    f"Customer base: {int(summary['total_customers']):,} active accounts across segments"
]

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

for idx, takeaway in enumerate(takeaways):
    if idx > 0:
        p = text_frame.add_paragraph()
    else:
        p = text_frame.paragraphs[0]
    
    p.text = takeaway
    p.font.size = Pt(20)
    p.level = 0
    p.space_before = Pt(12)

# >>> SAVE PRESENTATION
prs.save('output/financial_report.pptx')
print("\nSaved to: output/financial_report.pptx")

conn.close()

print("\n" + "=" * 70)
print("POWERPOINT GENERATION COMPLETE")
print("=" * 70)